from pptx import Presentation
from pptx.util import Inches

def create_presentation():
    # Taqdimot obyekti yaratish
    prs = Presentation()

    # Slaydlarni qo'shish
    titles = [
        "Iqtisodiy Axborotlarni Matn Redaktorlari Vositalaridan Foydalangan Holda Qayta Ishlash",
        "Microsoft Excel Dasturi - Ishchi Kitoblar",
        "Excel Jadvali va Kitoblari",
        "Jadval Yaratish va Tahrirlash",
        "Ma'lumotlarni Tahlil Qilish",
        "Moliyaviy Hisob-kitoblar",
        "Excel Diagrammalari va Grafikalar",
        "Prognozlash va Hisoblash",
        "Iqtisodiy Tahlilning Ahamiyati",
        "Xulosa"
    ]
    contents = [
        "Ushbu taqdimotda Microsoft Excel dasturining iqtisodiy axborotlarni qayta ishlashdagi roli va uning iqtisodiy tahlil uchun qanday ishlatilishini ko'rib chiqamiz.",
        "Excelda ishchi kitoblar bir nechta ishchi varaqlardan tashkil topgan bo'lib, ma'lumotlarni kiritish va tahlil qilishda yordam beradi.",
        "Excelda jadval yaratish va ularni tahrirlash juda oson. Jadvalda ustunlar va satrlar bo'lib, har bir ustun ma'lum turdagi ma'lumotni o'z ichiga oladi.",
        "Jadvalni yaratish uchun ma'lumotlarni kiritib, keyin `Insert > Table` opsiyasidan foydalanib, ular to'g'ri shaklda tashkil etiladi.",
        "Excelda iqtisodiy ma'lumotlarni tahlil qilish uchun turli statistik formulalar va diagrammalar ishlatiladi.",
        "Moliyaviy hisobotlar, balanslar va naqd pul oqimi Excelda aniq va oson hisoblanadi. Excel formulalari yordamida hisob-kitoblar avtomatlashtiriladi.",
        "Excelda diagramma va grafikalar yaratish, ma'lumotlarni vizual tarzda taqdim etish uchun juda qulay. Bu diagrammalar iqtisodiy qarorlar qabul qilishda yordam beradi.",
        "Excelda vaqt ketma-ketligi ma'lumotlarini tahlil qilish va prognozlash imkoniyatlari mavjud. Bu, ayniqsa, iqtisodiy ko'rsatkichlarni prognoz qilishda foydalidir.",
        "Iqtisodiy tahlil qilishda Excelning samaradorligi, ma'lumotlar bazasi va statistik vositalarining ishlatilishi muhim ahamiyatga ega.",
        "Iqtisodiy axborotlarni qayta ishlashda Excelni ishlatish orqali samarali va tez qarorlar qabul qilish imkoniyatini yaratadi."
    ]

    # Slaydlarni yaratish
    for i in range(len(titles)):
        slide_layout = prs.slide_layouts[1]  # Title and Content slayd dizayni
        slide = prs.slides.add_slide(slide_layout)

        # Sarlavha va matn qo'shish
        title = slide.shapes.title
        content = slide.shapes.placeholders[1]

        title.text = titles[i]
        content.text = contents[i]

    # Taqdimotni saqlash
    prs.save("Iqtisodiy_Axborot_Taqdimoti.pptx")
    print("Taqdimot yaratildi: Iqtisodiy_Axborot_Taqdimoti.pptx")

# Taqdimotni yaratish
if __name__ == "__main__":
    create_presentation()

